"""Collision-safe FlowGuard judge deck.

Design reminder: restrained operations-console aesthetic on a white canvas;
each slide uses a fixed header, content band, and footer to avoid imported
template placeholders reflowing or overlapping in Canva.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/ubuntu")
ASSETS = ROOT / "webdev-static-assets"
OUTPUT = ASSETS / "FlowGuard_CodeForge_Judge_Presentation_v4.pptx"

NAVY = "0B1730"
INK = "17223B"
MUTED = "64748B"
BLUE = "1F64E8"
PALE_BLUE = "EAF2FF"
TEAL = "0E9F8B"
AMBER = "E8A219"
RED = "CE4C4C"
LINE = "D9E1EC"
OFF_WHITE = "F8FAFC"
WHITE = "FFFFFF"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill=WHITE, line=None, rounded=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    if line:
        shape.line.color.rgb = color(line)
        shape.line.width = Pt(0.65)
    else:
        shape.line.fill.background()
    return shape


def text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=14,
    fill=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    font="Arial",
    spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    for index, line in enumerate(value.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = spacing
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color(fill)
    return box


def picture(slide, name, x, y, w, h):
    source = ASSETS / name
    if source.exists():
        slide.shapes.add_picture(str(source), Inches(x), Inches(y), Inches(w), Inches(h))


def rule(slide, x, y, w, fill=BLUE, height=0.035):
    return rect(slide, x, y, w, height, fill=fill)


def brand_header(slide, page, section):
    logo = ASSETS / "flowguard-mark.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.62), Inches(0.39), height=Inches(0.30))
    text(slide, 1.00, 0.42, 1.50, 0.16, "FLOWGUARD", size=10, fill=INK, bold=True)
    text(slide, 1.00, 0.59, 1.50, 0.12, "FORGEVI", size=6.4, fill=MUTED, bold=True)
    text(slide, 9.40, 0.42, 3.25, 0.14, "CODEFORGE 2026  •  JUDGE DECK", size=7.2, fill=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    rule(slide, 0.62, 0.91, 12.10, fill=LINE, height=0.012)
    text(slide, 0.64, 1.08, 2.00, 0.16, f"0{page}  /  {section}", size=7.5, fill=BLUE, bold=True)


def title(slide, page, section, heading, subtitle=None):
    brand_header(slide, page, section)
    text(slide, 0.64, 1.34, 11.70, 0.43, heading, size=26, fill=INK, bold=True)
    rule(slide, 0.64, 1.91, 1.08, fill=BLUE)
    if subtitle:
        text(slide, 0.64, 2.08, 11.4, 0.23, subtitle, size=10.5, fill=MUTED)


def footer(slide, page):
    rule(slide, 0.62, 7.05, 12.10, fill=LINE, height=0.012)
    text(slide, 0.64, 7.15, 5.0, 0.12, "FORGEVI  •  FLOWGUARD  •  DURABLE HUMAN WORKFLOWS", size=6.4, fill=MUTED, bold=True)
    text(slide, 12.28, 7.14, 0.30, 0.12, str(page), size=7, fill=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def panel(slide, x, y, w, h, accent=BLUE, fill=WHITE):
    rect(slide, x, y, w, h, fill=fill, line=LINE, rounded=True)
    rule(slide, x, y, w, fill=accent, height=0.045)


def cover(slide):
    rect(slide, 0, 0, 13.333, 3.18, fill="F5F1E9")
    text(slide, 0.62, 0.15, 4.50, 1.08, "CODE", size=55, fill="111111", bold=True, font="Georgia", spacing=0.88)
    text(slide, 0.62, 1.28, 4.90, 1.08, "FORGE", size=55, fill="111111", bold=True, font="Georgia", spacing=0.88)
    logo = ASSETS / "flowguard-mark.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(8.82), Inches(0.43), height=Inches(0.38))
    text(slide, 9.30, 0.48, 1.70, 0.15, "FLOWGUARD", size=10.5, fill=INK, bold=True)
    text(slide, 9.30, 0.67, 1.70, 0.12, "FORGEVI", size=6.5, fill=MUTED, bold=True)
    college_logo = ASSETS / "ghraisoni-skilltech-nagpur-logo.png"
    if college_logo.exists():
        slide.shapes.add_picture(str(college_logo), Inches(11.30), Inches(0.22), Inches(1.55), Inches(0.81))
    rule(slide, 0, 3.18, 13.333, fill=WHITE, height=0.02)
    text(slide, 0.66, 3.42, 3.05, 0.20, "Team Name:", size=13.5, fill=MUTED)
    text(slide, 0.66, 3.73, 4.55, 0.40, "ForgeVI", size=27, fill=INK, bold=True)
    text(slide, 7.26, 3.42, 2.65, 0.20, "Problem Code:", size=13.5, fill=MUTED)
    text(slide, 7.26, 3.73, 3.70, 0.33, "P-02", size=23, fill=INK, bold=True)
    text(slide, 0.66, 4.42, 4.05, 0.20, "Problem Statement Title:", size=13.5, fill=MUTED)
    text(slide, 0.66, 4.73, 11.76, 0.64, "Distributed Transaction Coordinator for Human Workflows", size=22, fill=INK, bold=True, spacing=1.02)
    text(slide, 0.66, 5.60, 2.35, 0.20, "Team Members:", size=13.5, fill=MUTED)
    text(slide, 0.66, 5.92, 11.42, 0.50, "Aditya Devhare — Product Presenter & Team Leader  •  Arpit Tagade — Backend Developer\nRohan Kodane — Tester  •  Atharva Andhare — Frontend Developer", size=11.4, fill=INK, spacing=1.06)
    text(slide, 0.66, 6.68, 11.55, 0.16, "FlowGuard retains human control while making every business hand-off durable, recoverable, and auditable.", size=9.8, fill=BLUE, bold=True)
    footer(slide, 1)


def problem(slide):
    title(slide, 1, "THE PROBLEM", "A purchase request can fail between systems", "CRM, inventory, payment, invoicing, notification, and a human manager do not share one transaction boundary.")
    text(slide, 0.66, 2.52, 6.75, 0.40, "A normal request-response flow can duplicate an action, lose progress after failure, or leave a manager decision outside the audit trail.", size=13.4, fill=INK, bold=True)
    points = [
        ("1", "Partial failure", "A later service fails after an earlier action succeeds.", RED),
        ("2", "Human delay", "A decision takes minutes, not one HTTP request.", AMBER),
        ("3", "Safe retry", "A retry must not create a duplicate business effect.", TEAL),
    ]
    for idx, (number, heading, body, accent) in enumerate(points):
        y = 3.18 + idx * 0.98
        panel(slide, 0.66, y, 6.75, 0.78, accent=accent)
        rect(slide, 0.90, y + 0.20, 0.35, 0.35, fill=accent, line=None, rounded=True)
        text(slide, 0.90, y + 0.27, 0.35, 0.11, number, size=7.2, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, 1.48, y + 0.16, 2.02, 0.15, heading, size=11.2, fill=INK, bold=True)
        text(slide, 1.48, y + 0.39, 5.55, 0.15, body, size=9.3, fill=MUTED)
    rect(slide, 8.07, 2.50, 4.55, 3.58, fill=NAVY, line=None, rounded=True)
    picture(slide, "flowguard-operations-hero.png", 8.09, 2.52, 4.51, 3.54)
    text(slide, 8.10, 6.25, 4.50, 0.30, "One coordinator makes each hand-off visible, recoverable, and safe to retry.", size=10.1, fill=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2)


def guarantees(slide):
    title(slide, 2, "PROBLEM DECOMPOSITION", "The coordinator must preserve five guarantees", "The design treats reliability as persisted, observable workflow state—not a best-effort sequence of API calls.")
    items = [
        ("01", "Independent participants", "Persist state before the next action.", BLUE),
        ("02", "Uncertain retry", "Reuse a stable idempotency key.", TEAL),
        ("03", "Human delay", "Store owner, deadline, and decision.", AMBER),
        ("04", "Permanent failure", "Compensate completed actions in reverse order.", RED),
        ("05", "Multiple workers", "Atomically lease one due job.", BLUE),
    ]
    locations = [(0.66, 2.62), (4.52, 2.62), (8.38, 2.62), (0.66, 4.20), (4.52, 4.20)]
    for (num, heading, body, accent), (x, y) in zip(items, locations):
        panel(slide, x, y, 3.25, 1.20, accent=accent)
        text(slide, x + 0.20, y + 0.23, 0.48, 0.16, num, size=8.0, fill=accent, bold=True)
        text(slide, x + 0.20, y + 0.49, 2.78, 0.17, heading, size=10.1, fill=INK, bold=True)
        text(slide, x + 0.20, y + 0.78, 2.78, 0.18, body, size=8.6, fill=MUTED)
    rect(slide, 8.38, 4.20, 4.24, 1.20, fill=PALE_BLUE, line=None, rounded=True)
    text(slide, 8.64, 4.45, 3.72, 0.14, "TECHNICAL CLAIM", size=7.6, fill=BLUE, bold=True)
    text(slide, 8.64, 4.70, 3.62, 0.42, "Durable execution + atomic claims + idempotency + events make human approval recoverable.", size=10.0, fill=INK, bold=True, spacing=1.05)
    footer(slide, 3)


def research(slide):
    title(slide, 3, "RESEARCH & PRIOR ART", "Research favors orchestration with compensation", "Each design choice answers a known distributed-systems failure mode and is visible in the prototype.")
    cards = [
        ("Saga pattern", "Sequence local transactions; compensate when a later action fails.", BLUE),
        ("Compensating commands", "Undo actions must be idempotent and may require human recovery.", TEAL),
        ("Atomic job claim", "One MongoDB compound update leases one due job safely.", AMBER),
        ("Durable human task", "Store owner, deadline, and decision—not an open browser request.", RED),
    ]
    for idx, (heading, body, accent) in enumerate(cards):
        x = 0.66 + (idx % 2) * 6.08
        y = 2.60 + (idx // 2) * 1.64
        panel(slide, x, y, 5.42, 1.32, accent=accent)
        text(slide, x + 0.24, y + 0.27, 4.82, 0.16, heading, size=12.2, fill=INK, bold=True)
        text(slide, x + 0.24, y + 0.61, 4.84, 0.38, body, size=10.0, fill=MUTED, spacing=1.08)
    rect(slide, 0.66, 5.68, 11.50, 0.60, fill=PALE_BLUE, line=None, rounded=True)
    text(slide, 0.94, 5.90, 10.88, 0.16, "FLOWGUARD GAP  •  Retry, human wait, rejection, and recovery are visible in one judgeable system.", size=10.3, fill=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4)


def approach(slide):
    title(slide, 4, "PROPOSED APPROACH", "FlowGuard is a persisted seven-step Saga", "Every action writes durable workflow state; rejection and permanent failure compensate completed actions in reverse business order.")
    steps = [
        ("1", "Create CRM\norder", BLUE),
        ("2", "Reserve\ninventory", BLUE),
        ("3", "Authorize\npayment", BLUE),
        ("4", "Manager\napproval", AMBER),
        ("5", "Capture\npayment", TEAL),
        ("6", "Create\ninvoice", TEAL),
        ("7", "Send\nnotification", TEAL),
    ]
    slots = [(0.76, 2.72), (3.53, 2.72), (6.30, 2.72), (9.07, 2.72), (2.15, 4.18), (5.00, 4.18), (7.85, 4.18)]
    for (number, label, accent), (x, y) in zip(steps, slots):
        panel(slide, x, y, 1.85, 0.92, accent=accent)
        rect(slide, x + 0.16, y + 0.22, 0.31, 0.31, fill=accent, line=None, rounded=True)
        text(slide, x + 0.16, y + 0.29, 0.31, 0.11, number, size=7, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, x + 0.58, y + 0.22, 1.03, 0.42, label, size=8.9, fill=INK, bold=True, spacing=0.94)
    text(slide, 11.09, 3.39, 0.45, 0.18, "↓", size=20, fill=BLUE, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0.76, 5.52, 11.42, 0.53, fill=NAVY, line=None, rounded=True)
    text(slide, 1.02, 5.70, 10.92, 0.15, "The worker resumes only the next safe action; the human decision remains durable state.", size=10.4, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 5)


def architecture(slide):
    title(slide, 5, "SYSTEM ARCHITECTURE", "Durable state separates the UI from the work", "The interface submits and explains; the API and lease-safe worker coordinate durable state and recovery.")
    systems = [
        (0.66, "React operations\nconsole", "Vercel\nHTTPS + JWT", BLUE),
        (4.35, "Express API +\nlease-safe worker", "Render\ncommand + recovery", TEAL),
        (8.04, "MongoDB Atlas", "executions • jobs\napprovals • events • users", AMBER),
    ]
    for idx, (x, heading, body, accent) in enumerate(systems):
        panel(slide, x, 2.68, 3.05, 1.36, accent=accent)
        text(slide, x + 0.22, 2.99, 2.60, 0.35, heading, size=11.5, fill=INK, bold=True, align=PP_ALIGN.CENTER, spacing=0.98)
        text(slide, x + 0.22, 3.53, 2.60, 0.27, body, size=8.4, fill=MUTED, align=PP_ALIGN.CENTER, spacing=1.00)
        if idx < 2:
            text(slide, x + 3.15, 3.20, 0.52, 0.18, "→", size=18, fill=BLUE, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 0.66, 4.64, 6.45, 1.12, fill=PALE_BLUE, line=None, rounded=True)
    text(slide, 0.95, 4.90, 5.82, 0.14, "ROLE BOUNDARIES", size=7.8, fill=BLUE, bold=True)
    text(slide, 0.95, 5.18, 5.72, 0.32, "Requester submits own work  •  Operator observes  •  Manager decides\nAdministrator recovers under authorization and audits every action", size=9.3, fill=INK, bold=True, spacing=1.05)
    rect(slide, 7.55, 4.64, 4.61, 1.12, fill=NAVY, line=None, rounded=True)
    picture(slide, "flowguard-product-hero.png", 7.57, 4.66, 4.57, 1.08)
    footer(slide, 6)


def impact(slide):
    title(slide, 6, "INNOVATION & IMPACT", "A durable human decision makes recovery demonstrable", "The prototype establishes evidence for correctness, resilience, authorization, and an operable human-in-the-loop path.")
    metrics = [("6", "focused engine / seed tests", BLUE), ("14", "end-to-end API smoke scenarios", TEAL), ("5 min", "live approval decision window", AMBER)]
    for idx, (metric, label, accent) in enumerate(metrics):
        x = 0.66 + idx * 2.22
        rect(slide, x, 2.60, 1.92, 1.18, fill=PALE_BLUE, line=None, rounded=True)
        text(slide, x, 2.87, 1.92, 0.22, metric, size=19, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        text(slide, x + 0.12, 3.27, 1.68, 0.18, label, size=7.3, fill=MUTED, bold=True, align=PP_ALIGN.CENTER)
    panel(slide, 0.66, 4.22, 5.92, 1.28, accent=BLUE)
    text(slide, 0.93, 4.49, 5.34, 0.14, "EVIDENCE", size=7.7, fill=BLUE, bold=True)
    text(slide, 0.93, 4.74, 5.34, 0.22, "Approval completion • compensation • idempotency • audit trail", size=11.2, fill=INK, bold=True)
    text(slide, 0.93, 5.11, 5.34, 0.15, "A fresh live workflow reached an OPEN approval task before expiry.", size=9.0, fill=MUTED)
    rect(slide, 7.07, 2.60, 2.45, 2.90, fill=NAVY, line=None, rounded=True)
    picture(slide, "flowguard-approval-evidence-v2.png", 7.09, 2.62, 2.41, 2.86)
    rect(slide, 9.73, 2.60, 2.45, 2.90, fill=NAVY, line=None, rounded=True)
    picture(slide, "flowguard-recovery-evidence-v2.png", 9.75, 2.62, 2.41, 2.86)
    text(slide, 0.66, 5.90, 11.50, 0.22, "REAL-WORLD VALUE  •  Procurement, fulfilment, claims, and operations retain human judgment without losing process continuity.", size=10.3, fill=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 7)


def defense(slide):
    title(slide, 7, "REFERENCES & DEFENSE", "Evidence, scope, and technical defense", "FlowGuard claims a narrow, demonstrable coordinator—not real settlement or a generic workflow builder.")
    panel(slide, 0.66, 2.60, 7.05, 3.02, accent=BLUE)
    text(slide, 0.94, 2.89, 5.80, 0.14, "REFERENCES", size=7.7, fill=BLUE, bold=True)
    references = "1. Azure Architecture Center — Saga distributed transactions pattern\n2. Azure Architecture Center — Compensating Transaction pattern\n3. MongoDB Documentation — Compound operations\n4. Temporal Documentation — Human-in-the-loop workflow pattern\n5. CODEFORGE 2026 — Participant Evaluation Framework"
    text(slide, 0.94, 3.22, 6.15, 1.55, references, size=9.4, fill=INK, spacing=1.18)
    rect(slide, 8.02, 2.60, 4.16, 3.02, fill=PALE_BLUE, line=None, rounded=True)
    text(slide, 8.34, 2.90, 3.52, 0.14, "PROTOTYPE BOUNDARY", size=7.7, fill=BLUE, bold=True)
    text(slide, 8.34, 3.22, 3.48, 0.50, "One predefined purchase workflow with deterministic adapters.", size=11.4, fill=INK, bold=True, spacing=1.04)
    text(slide, 8.34, 3.96, 3.46, 0.72, "No claim of real settlement, generic workflow building, or multi-stage voting.\n\nAI assisted design, code scaffolding, debugging, documentation, and testing; ForgeVI owns the evidence and technical defense.", size=8.4, fill=MUTED, spacing=1.06)
    rect(slide, 0.66, 5.93, 11.52, 0.48, fill=NAVY, line=None, rounded=True)
    text(slide, 0.95, 6.09, 10.92, 0.15, "Judge demo: requester submits  →  manager decides  →  evidence shows completion, retry, or compensation.", size=9.6, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8)


def build():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    functions = [cover, problem, guarantees, research, approach, architecture, impact, defense]
    for function in functions:
        slide = presentation.slides.add_slide(blank)
        rect(slide, 0, 0, 13.333, 7.5, fill=WHITE)
        function(slide)
    presentation.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
