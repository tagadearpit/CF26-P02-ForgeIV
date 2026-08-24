from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path('/home/ubuntu')
TEMPLATE = ROOT / 'upload' / 'CodeForge_PPT_Template.pptx'
ASSETS = ROOT / 'webdev-static-assets'
OUTPUT = ASSETS / 'FlowGuard_CodeForge_Judge_Presentation.pptx'

NAVY = '0E1B36'
INK = '16233D'
MUTED = '667085'
BLUE = '2563EB'
BLUE_SOFT = 'EAF1FF'
TEAL = '0F9E8A'
AMBER = 'F59E0B'
RED = 'D14343'
LINE = 'D7DEE8'
WHITE = 'FFFFFF'


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def clear_text(shape):
    if shape.has_text_frame:
        shape.text_frame.clear()


def clear_slide_text(slide, keep_codeforge=False):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip().upper()
            if keep_codeforge and ('CODE' in text or 'FORGE' in text):
                continue
            clear_text(shape)


def set_fill(shape, color, transparency=0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)
    fill.transparency = transparency


def set_line(shape, color=LINE, width=0.6):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill, transparency)
    if line:
        set_line(shape, line)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
             font='Arial', valign=MSO_ANCHOR.TOP, margin=0.04, line_spacing=1.12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    paragraphs = text.split('\n')
    for index, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, x, y, w, h, lines, fill=None, line=None, radius=True):
    if fill:
        add_rect(slide, x, y, w, h, fill=fill, line=line, radius=radius)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.11)
    tf.margin_bottom = Inches(0.10)
    for index, spec in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = spec.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(spec.get('after', 2))
        p.line_spacing = spec.get('line_spacing', 1.1)
        run = p.add_run()
        run.text = spec['text']
        run.font.name = 'Arial'
        run.font.size = Pt(spec.get('size', 14))
        run.font.bold = spec.get('bold', False)
        run.font.color.rgb = rgb(spec.get('color', INK))
    return box


def add_rule(slide, x, y, w, color=BLUE, height=0.035):
    return add_rect(slide, x, y, w, height, fill=color, radius=False)


def add_title(slide, number, title, subtitle=None):
    add_text(slide, 0.34, 0.25, 2.0, 0.2, f'0{number}  /  FLOWGUARD', size=7.8, color=BLUE, bold=True)
    add_text(slide, 0.34, 0.48, 8.6, 0.48, title, size=24, color=INK, bold=True)
    add_rule(slide, 0.34, 1.06, 1.05, color=BLUE)
    if subtitle:
        add_text(slide, 0.34, 1.16, 8.9, 0.26, subtitle, size=10.5, color=MUTED)


def add_footer(slide, index):
    add_rule(slide, 0.34, 5.30, 9.32, color=LINE, height=0.012)
    add_text(slide, 0.34, 5.36, 5.5, 0.15, 'FORGEVI  •  CODEFORGE 2026  •  FLOWGUARD', size=6.6, color=MUTED, bold=True)
    add_text(slide, 9.18, 5.34, 0.46, 0.16, str(index), size=7, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_logo(slide, x, y, size=0.25):
    logo = ASSETS / 'flowguard-mark.png'
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(x), Inches(y), height=Inches(size))


def add_step(slide, x, y, number, title, body, accent=BLUE):
    add_rect(slide, x, y, 1.78, 1.23, fill=WHITE, line=LINE)
    add_rect(slide, x + 0.16, y + 0.14, 0.30, 0.30, fill=accent, radius=True)
    add_text(slide, x + 0.16, y + 0.19, 0.30, 0.17, str(number), size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.16, y + 0.52, 1.45, 0.24, title, size=10.5, color=INK, bold=True)
    add_text(slide, x + 0.16, y + 0.82, 1.46, 0.29, body, size=7.6, color=MUTED)


def cover(slide):
    clear_slide_text(slide, keep_codeforge=True)
    add_logo(slide, 5.54, 0.36, 0.30)
    add_text(slide, 5.92, 0.40, 2.0, 0.18, 'FLOWGUARD', size=10, color=INK, bold=True)
    add_text(slide, 5.92, 0.59, 2.0, 0.13, 'FORGEVI', size=6.5, color=MUTED, bold=True)
    hero = ASSETS / 'flowguard-product-hero.png'
    if hero.exists():
        slide.shapes.add_picture(str(hero), Inches(5.52), Inches(0.93), Inches(4.11), Inches(2.45))
        add_rect(slide, 5.52, 2.78, 4.11, 0.60, fill=NAVY, transparency=15, radius=False)
        add_text(slide, 5.74, 2.92, 3.60, 0.22, 'DECISION DUE  •  HUMAN CONTROL RETAINED', size=8.5, color=WHITE, bold=True)
    add_text(slide, 5.54, 3.63, 3.86, 0.46, 'P-02 — Distributed Transaction Coordinator for Human Workflows', size=17, color=INK, bold=True)
    add_rule(slide, 5.54, 4.20, 1.08, color=BLUE)
    add_text(slide, 5.54, 4.38, 3.80, 0.31, 'A durable coordinator for human-approved business workflows.', size=10.5, color=MUTED)
    add_rich_text(slide, 5.54, 4.82, 3.82, 0.48, [
        {'text': 'ForgeVI  •  Aditya Devhare  •  Arpit Tagade  •  Rohan Kodane  •  Atharva Andhare', 'size': 7.8, 'color': INK, 'bold': True, 'line_spacing': 1.1}
    ], fill=BLUE_SOFT, line=None)


def problem(slide):
    clear_slide_text(slide)
    add_title(slide, 1, 'A purchase request can fail between systems', 'THE PROBLEM')
    add_text(slide, 0.46, 1.62, 4.18, 0.72, 'CRM, inventory, payment, invoicing, notification, and a human manager do not share one transaction boundary.', size=17, color=INK, bold=True)
    add_text(slide, 0.46, 2.45, 4.10, 0.58, 'A normal request-response flow can duplicate an action, lose progress after a failure, or leave a manager decision outside the audit trail.', size=11.4, color=MUTED)
    add_step(slide, 0.46, 3.40, 1, 'Partial failure', 'A later service fails after an earlier step succeeds.', RED)
    add_step(slide, 2.30, 3.40, 2, 'Human delay', 'A decision takes minutes, not one HTTP request.', AMBER)
    add_step(slide, 4.14, 3.40, 3, 'Safe retry', 'A retry must not create a duplicate effect.', TEAL)
    image = ASSETS / 'flowguard-operations-hero.png'
    if image.exists():
        slide.shapes.add_picture(str(image), Inches(6.10), Inches(1.40), Inches(3.28), Inches(3.50))
    add_footer(slide, 2)


def decomposition(slide):
    clear_slide_text(slide)
    add_title(slide, 2, 'The coordinator must preserve five guarantees', 'PROBLEM DECOMPOSITION')
    entries = [
        ('Independent participants', 'Persist state before the next action.'),
        ('Uncertain retry', 'Reuse a stable idempotency key.'),
        ('Human delay', 'Store owner, deadline, and decision.'),
        ('Permanent failure', 'Compensate completed actions in reverse order.'),
        ('Multiple workers', 'Atomically lease one due job.'),
    ]
    for idx, (head, body) in enumerate(entries):
        y = 1.55 + idx * 0.64
        add_rect(slide, 0.48, y, 0.26, 0.26, fill=BLUE if idx < 3 else TEAL, radius=True)
        add_text(slide, 0.48, y + 0.045, 0.26, 0.13, str(idx + 1), size=7.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, 0.90, y - 0.015, 2.50, 0.19, head, size=11, color=INK, bold=True)
        add_text(slide, 3.48, y - 0.010, 4.35, 0.24, body, size=10.5, color=MUTED)
    add_rich_text(slide, 0.48, 4.64, 8.92, 0.46, [
        {'text': 'TECHNICAL CLAIM  •  Persisted execution state + atomic job claims + stable idempotency + append-only events make delayed human approval recoverable.', 'size': 9, 'color': INK, 'bold': True}
    ], fill=BLUE_SOFT)
    add_footer(slide, 3)


def research(slide):
    clear_slide_text(slide)
    add_title(slide, 3, 'Research favors orchestration with compensation', 'RESEARCH & PRIOR ART')
    cards = [
        ('Saga pattern', 'Sequence local transactions; compensate when a later action fails.', BLUE),
        ('Compensating commands', 'Undo actions must be idempotent and may need human recovery.', TEAL),
        ('Atomic job claim', 'One MongoDB compound update leases one due job.', AMBER),
        ('Durable human task', 'Store owner, deadline, and decision—not an open browser request.', RED),
    ]
    for idx, (head, body, accent) in enumerate(cards):
        col = idx % 2
        row = idx // 2
        x = 0.48 + col * 4.52
        y = 1.55 + row * 1.48
        add_rect(slide, x, y, 4.06, 1.18, fill=WHITE, line=LINE)
        add_rule(slide, x, y, 0.86, color=accent)
        add_text(slide, x + 0.20, y + 0.24, 3.52, 0.24, head, size=12, color=INK, bold=True)
        add_text(slide, x + 0.20, y + 0.58, 3.54, 0.38, body, size=9.6, color=MUTED)
    add_rich_text(slide, 0.48, 4.63, 8.92, 0.44, [
        {'text': 'FLOWGUARD GAP  •  The prototype makes retry, human wait, rejection, and recovery visible in one judgeable system.', 'size': 9.4, 'color': INK, 'bold': True}
    ], fill=BLUE_SOFT)
    add_footer(slide, 4)


def approach(slide):
    clear_slide_text(slide)
    add_title(slide, 4, 'FlowGuard is a persisted seven-step Saga', 'PROPOSED APPROACH')
    steps = [
        ('1', 'Create\nCRM order', BLUE),
        ('2', 'Reserve\ninventory', BLUE),
        ('3', 'Authorize\npayment', BLUE),
        ('4', 'Manager\napproval', AMBER),
        ('5', 'Capture\npayment', TEAL),
        ('6', 'Create\ninvoice', TEAL),
        ('7', 'Send\nnotification', TEAL),
    ]
    start_x = 0.47
    for idx, (num, label, accent) in enumerate(steps):
        x = start_x + idx * 1.29
        add_rect(slide, x, 2.08, 1.03, 0.91, fill=WHITE, line=accent)
        add_rect(slide, x + 0.09, 2.17, 0.23, 0.23, fill=accent, radius=True)
        add_text(slide, x + 0.09, 2.205, 0.23, 0.10, num, size=6.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.10, 2.49, 0.80, 0.27, label, size=8.1, color=INK, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_rule(slide, x + 1.04, 2.51, 0.21, color=LINE, height=0.022)
    add_rich_text(slide, 1.10, 3.48, 7.95, 0.67, [
        {'text': 'FORWARD PATH  •  Every action creates durable execution, step, job, idempotency, and event records.', 'size': 10.2, 'color': INK, 'bold': True},
        {'text': 'REJECTION / PERMANENT FAILURE  •  Compensate completed actions in reverse business order.', 'size': 9.4, 'color': RED, 'bold': True},
    ], fill=BLUE_SOFT)
    add_text(slide, 1.10, 4.40, 7.95, 0.25, 'The worker resumes only the next safe action; a human decision remains durable state.', size=10.6, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)


def architecture(slide):
    clear_slide_text(slide)
    add_title(slide, 5, 'Durable state separates the UI from the work', 'SYSTEM ARCHITECTURE')
    columns = [
        (0.50, 'React operations\nconsole', 'Vercel\nHTTPS + JWT', BLUE),
        (3.26, 'Express API +\nlease-safe worker', 'Render\ncommand + recovery', TEAL),
        (6.02, 'MongoDB Atlas', 'executions • jobs\napprovals • events • users', AMBER),
    ]
    for idx, (x, head, body, accent) in enumerate(columns):
        add_rect(slide, x, 1.72, 2.10, 1.22, fill=WHITE, line=LINE)
        add_rule(slide, x, 1.72, 2.10, color=accent)
        add_text(slide, x + 0.16, 2.02, 1.78, 0.35, head, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 2.50, 1.74, 0.24, body, size=8.6, color=MUTED, align=PP_ALIGN.CENTER)
        if idx < 2:
            add_text(slide, x + 2.20, 2.22, 0.36, 0.24, '→', size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rich_text(slide, 0.52, 3.40, 5.35, 1.20, [
        {'text': 'ROLE BOUNDARIES', 'size': 8, 'color': BLUE, 'bold': True},
        {'text': 'Requester  →  submits and sees own work', 'size': 10.2, 'color': INK, 'bold': True},
        {'text': 'Operator  →  observes • Manager  →  decides • Administrator  →  recovers and audits', 'size': 9.1, 'color': MUTED},
    ], fill=BLUE_SOFT)
    screenshot = ASSETS / 'flowguard-product-hero.png'
    if screenshot.exists():
        slide.shapes.add_picture(str(screenshot), Inches(6.35), Inches(3.35), Inches(3.10), Inches(1.72))
    add_footer(slide, 6)


def impact(slide):
    clear_slide_text(slide)
    add_title(slide, 6, 'A durable human decision makes recovery demonstrable', 'INNOVATION & IMPACT')
    metrics = [
        ('6', 'focused engine / seed tests'),
        ('14', 'end-to-end API smoke scenarios'),
        ('5 min', 'live approval decision window'),
    ]
    for idx, (metric, label) in enumerate(metrics):
        x = 0.50 + idx * 1.72
        add_rect(slide, x, 1.58, 1.47, 0.93, fill=BLUE_SOFT, line=None)
        add_text(slide, x, 1.74, 1.47, 0.22, metric, size=19, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.08, 2.10, 1.30, 0.18, label, size=7.0, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rich_text(slide, 0.50, 2.82, 4.94, 1.33, [
        {'text': 'EVIDENCE', 'size': 8, 'color': BLUE, 'bold': True},
        {'text': 'Approval completion • compensation • idempotency • recovery authorization • audit trail', 'size': 12.5, 'color': INK, 'bold': True, 'line_spacing': 1.12},
        {'text': 'A fresh live workflow reached an OPEN approval task before expiry.', 'size': 9.3, 'color': MUTED},
    ], fill=WHITE, line=LINE)
    approval = ASSETS / 'flowguard-human-approval.png'
    recovery = ASSETS / 'flowguard-safe-recovery.png'
    if approval.exists():
        slide.shapes.add_picture(str(approval), Inches(5.83), Inches(1.54), Inches(1.72), Inches(2.65))
    if recovery.exists():
        slide.shapes.add_picture(str(recovery), Inches(7.72), Inches(1.54), Inches(1.72), Inches(2.65))
    add_text(slide, 0.50, 4.55, 8.90, 0.32, 'REAL-WORLD VALUE  •  Procurement, fulfilment, claims, and operations retain human judgment without losing process continuity.', size=10.4, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)


def references(slide):
    clear_slide_text(slide)
    add_title(slide, 7, 'Evidence, scope, and technical defense', 'REFERENCES & DEFENSE')
    add_rich_text(slide, 0.50, 1.50, 5.64, 2.55, [
        {'text': 'REFERENCES', 'size': 8, 'color': BLUE, 'bold': True},
        {'text': '1. Azure Architecture Center — Saga distributed transactions pattern', 'size': 9.2, 'color': INK},
        {'text': '2. Azure Architecture Center — Compensating Transaction pattern', 'size': 9.2, 'color': INK},
        {'text': '3. MongoDB Documentation — Compound operations', 'size': 9.2, 'color': INK},
        {'text': '4. Temporal Documentation — Human-in-the-loop workflow pattern', 'size': 9.2, 'color': INK},
        {'text': '5. CODEFORGE 2026 — Participant Evaluation Framework', 'size': 9.2, 'color': INK},
    ], fill=WHITE, line=LINE)
    add_rich_text(slide, 6.40, 1.50, 3.02, 2.55, [
        {'text': 'PROTOTYPE BOUNDARY', 'size': 8, 'color': BLUE, 'bold': True},
        {'text': 'One predefined purchase workflow with deterministic adapters.', 'size': 11.3, 'color': INK, 'bold': True},
        {'text': 'No claim of real settlement, a generic workflow builder, or multi-stage voting.', 'size': 9.0, 'color': MUTED},
        {'text': 'AI disclosure: assisted design, code scaffolding, debugging, documentation, and testing. ForgeVI owns all claims and evidence.', 'size': 8.4, 'color': MUTED},
    ], fill=BLUE_SOFT, line=None)
    add_rect(slide, 0.50, 4.45, 8.92, 0.47, fill=NAVY, radius=True)
    add_text(slide, 0.76, 4.57, 8.40, 0.16, 'Judge demo: requester submits  →  manager decides  →  evidence shows completion, retry, or compensation.', size=9.3, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)


def main():
    prs = Presentation(str(TEMPLATE))
    if len(prs.slides) != 8:
        raise RuntimeError('Expected an eight-slide CodeForge template.')
    cover(prs.slides[0])
    problem(prs.slides[1])
    decomposition(prs.slides[2])
    research(prs.slides[3])
    approach(prs.slides[4])
    architecture(prs.slides[5])
    impact(prs.slides[6])
    references(prs.slides[7])
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == '__main__':
    main()
